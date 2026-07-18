# -*- coding: utf-8 -*-
"""Build the Stage 10 group-meeting PowerPoint deck.

The deck is intentionally generated from the existing Stage 10 meeting figure
package.  It does not rerun upstream science stages and does not download data.
"""
from __future__ import annotations

import argparse
import csv
import json
import re
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path


def build_parser() -> argparse.ArgumentParser:
    return argparse.ArgumentParser(description=__doc__)


if __name__ == "__main__" and any(arg in {"-h", "--help"} for arg in sys.argv[1:]):
    build_parser().parse_args()


try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - environment diagnostic
    raise SystemExit(
        "python-pptx is required. Use the bundled Codex Python runtime or install python-pptx."
    ) from exc

SCRIPT_PATH = Path(__file__).resolve()
CODE_ROOT = SCRIPT_PATH.parents[1]
if str(CODE_ROOT) not in sys.path:
    sys.path.insert(0, str(CODE_ROOT))

import path_config  # noqa: E402


STAGE_ID = "stage_10"
PROJECT_STAGE_ID = "geo_ring_cloud.stage_10"
RUN_ID = "stage_10_meeting_figures_202403"
DECK_ID = "stage_10_group_meeting_progress_presentation_cn_audience"

RUN_ROOT = path_config.RUNS_ROOT / RUN_ID
FIGURE_DIR = RUN_ROOT / "figures"
SOURCE_DIR = RUN_ROOT / "source_data"
REPORT_DIR = RUN_ROOT / "reports"
LOG_DIR = RUN_ROOT / "logs"
PPTX_DIR = RUN_ROOT / "pptx"

PLOT_INDEX = LOG_DIR / "stage_10_group_meeting_plot_index.csv"
FIGURE_GUIDE = REPORT_DIR / "stage_10_group_meeting_figure_guide_cn.md"
INTERROGATION_GUIDE = REPORT_DIR / "stage_10_group_meeting_interrogation_guide_cn.md"

OUTPUT_PPTX = PPTX_DIR / f"{DECK_ID}.pptx"
RENDER_CHECK_PDF = PPTX_DIR / f"{DECK_ID}_render_check.pdf"
SLIDE_INDEX = LOG_DIR / "stage_10_group_meeting_ppt_slide_index.csv"
SPEAKER_NOTES = REPORT_DIR / "stage_10_group_meeting_ppt_speaker_notes_cn.md"
MANIFEST = LOG_DIR / "stage_10_group_meeting_ppt_manifest.json"
QA_REPORT = LOG_DIR / "stage_10_group_meeting_ppt_qa_report.json"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(24, 34, 48),
    "muted": RGBColor(94, 105, 119),
    "light": RGBColor(246, 248, 250),
    "line": RGBColor(213, 219, 226),
    "navy": RGBColor(33, 70, 111),
    "teal": RGBColor(39, 137, 132),
    "gold": RGBColor(196, 132, 45),
    "red": RGBColor(184, 79, 73),
    "green": RGBColor(81, 137, 84),
    "white": RGBColor(255, 255, 255),
}

FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"


def read_csv_rows(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return list(csv.DictReader(handle))


def as_float(value: str | None, default: float | None = None) -> float | None:
    if value is None or value == "":
        return default
    try:
        return float(value)
    except ValueError:
        return default


def fmt_num(value: float | int | None, ndigits: int = 3) -> str:
    if value is None:
        return "NA"
    if isinstance(value, int) or float(value).is_integer():
        return f"{int(value):,}"
    return f"{float(value):.{ndigits}f}"


def fmt_delta(value: float | None, unit: str = "") -> str:
    if value is None:
        return "NA"
    sign = "+" if value >= 0 else ""
    return f"{sign}{value:.3f}{unit}"


def load_plot_index() -> dict[str, dict[str, str]]:
    rows = read_csv_rows(PLOT_INDEX)
    return {row["figure_id"].split("_", 1)[0]: row for row in rows}


def load_key_metrics() -> dict[str, float | int | str]:
    fig02 = read_csv_rows(SOURCE_DIR / "stage_10_group_meeting_fig02_source.csv")
    d1 = next(row for row in fig02 if row["domain"] == "D1_both_cloud")
    d7 = next(row for row in fig02 if row["domain"] == "D7_high_cloud")

    fig05 = read_csv_rows(SOURCE_DIR / "stage_10_group_meeting_fig05_source.csv")
    regret_all = next(
        row
        for row in fig05
        if row["panel"] == "selection_regret" and row["pixel_group"] == "ALL_VALID_CTH"
    )
    regret_high = next(
        row
        for row in fig05
        if row["panel"] == "selection_regret" and row["pixel_group"] == "high_cloud"
    )
    clean = next(
        row
        for row in fig05
        if row["panel"] == "clean_boundary_metrics" and row["clean_group"] == "clean_core"
    )
    boundary = next(
        row
        for row in fig05
        if row["panel"] == "clean_boundary_metrics" and row["clean_group"] == "boundary_or_broken_cloud"
    )

    fig06 = read_csv_rows(SOURCE_DIR / "stage_10_group_meeting_fig06_source.csv")
    comp = {row["metric"]: as_float(row["value"]) for row in fig06 if row["panel"] == "composite_inventory"}
    box7_cth = next(
        row
        for row in fig06
        if row["panel"] == "cth_delta_d1" and row["method_name"] == "box_7x7"
    )
    box7_high = next(
        row
        for row in fig06
        if row["panel"] == "high_cloud_delta" and row["method_name"] == "box_7x7"
    )
    box7_cloud = next(
        row
        for row in fig06
        if row["panel"] == "cloud_mask_delta" and row["method_name"] == "box_7x7"
    )

    fig04 = read_csv_rows(SOURCE_DIR / "stage_10_group_meeting_fig04_source.csv")
    d1_band = next(
        row
        for row in fig04
        if row.get("panel") == "a_vs_b_sensitivity"
        and row.get("domain") == "D1_both_cloud"
        and row.get("group") == "ALL"
        and row.get("sensitivity_type") == "reference_band_metric_delta"
    )

    return {
        "d1_n": int(float(d1["n_valid_cth"])),
        "d1_bias": as_float(d1["bias_km"]),
        "d1_mae": as_float(d1["mae_km"]),
        "d1_rmse": as_float(d1["rmse_km"]),
        "d1_within2": as_float(d1["within_2km_fraction"]),
        "d7_mae": as_float(d7["mae_km"]),
        "regret_current": as_float(regret_all["current_selected_mae_km"]),
        "regret_oracle": as_float(regret_all["best_available_mae_km"]),
        "regret_delta": as_float(regret_all["selection_regret_mae_km"]),
        "regret_best_fraction": as_float(regret_all["selected_is_best_fraction"]),
        "regret_high": as_float(regret_high["selection_regret_mae_km"]),
        "clean_mae": as_float(clean["mae_km"]),
        "boundary_mae": as_float(boundary["mae_km"]),
        "composite_files": int(comp.get("open_ok_files") or 0),
        "explicit_psf_kernel": int(comp.get("explicit_psf_kernel_candidates") or 0),
        "explicit_weight": int(comp.get("explicit_weight_candidates") or 0),
        "box7_cth_mae_delta": as_float(box7_cth["delta_mae_km_vs_nearest"]),
        "box7_high_mae_delta": as_float(box7_high["delta_mae_km_vs_nearest"]),
        "box7_agree_delta": as_float(box7_cloud["delta_agreement_vs_nearest"]),
        "ab_d1_mae_delta": as_float(d1_band["b_minus_a_mae_km"]),
    }


def add_rect(slide, x, y, w, h, color, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=None, bold=False, align=None, font=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align or PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = font or FONT_CN
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["ink"]
    return box


def add_multiline(slide, x, y, w, h, lines, size=15, color=None, bullet=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color or COLORS["ink"]
        p.line_spacing = 1.05
        p.space_after = Pt(4)
        if bullet:
            p.level = 0
            p._p.get_or_add_pPr().set("marL", "260000")
            p._p.get_or_add_pPr().set("indent", "-180000")
    return box


def add_header(slide, section: str, title: str):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.18), COLORS["navy"])
    add_text(slide, Inches(0.45), Inches(0.34), Inches(1.85), Inches(0.24), section, 8.5, COLORS["teal"], True, font=FONT_EN)
    add_text(slide, Inches(0.45), Inches(0.56), Inches(10.9), Inches(0.42), title, 22, COLORS["ink"], True)
    add_text(
        slide,
        Inches(11.55),
        Inches(0.58),
        Inches(1.25),
        Inches(0.24),
        "Stage 10",
        8,
        COLORS["muted"],
        True,
        PP_ALIGN.RIGHT,
        FONT_EN,
    )


def add_footer(slide, slide_no: int, note: str = "Reference: EPIC A-band Effective Cloud Height; units in km unless noted."):
    add_rect(slide, Inches(0.45), Inches(7.13), Inches(12.0), Inches(0.01), COLORS["line"])
    add_text(slide, Inches(0.45), Inches(7.18), Inches(10.8), Inches(0.2), note, 7.5, COLORS["muted"])
    add_text(slide, Inches(12.05), Inches(7.18), Inches(0.75), Inches(0.2), f"{slide_no:02d}", 7.5, COLORS["muted"], False, PP_ALIGN.RIGHT, FONT_EN)


def add_metric_card(slide, x, y, w, h, label, value, sub="", accent="teal"):
    add_rect(slide, x, y, w, h, COLORS["white"], COLORS["line"], radius=True)
    add_rect(slide, x, y, Inches(0.08), h, COLORS[accent], radius=False)
    add_text(slide, x + Inches(0.18), y + Inches(0.13), w - Inches(0.28), Inches(0.22), label, 9, COLORS["muted"], True)
    add_text(slide, x + Inches(0.18), y + Inches(0.42), w - Inches(0.28), Inches(0.38), value, 18, COLORS["ink"], True, font=FONT_EN)
    if sub:
        add_text(slide, x + Inches(0.18), y + Inches(0.88), w - Inches(0.28), Inches(0.24), sub, 8, COLORS["muted"])


def add_picture_fit(slide, path: Path, x, y, w, h):
    if not path.exists():
        raise FileNotFoundError(path)
    pic = slide.shapes.add_picture(str(path), x, y, width=w)
    if pic.height > h:
        slide.shapes._spTree.remove(pic._element)
        pic = slide.shapes.add_picture(str(path), x, y, height=h)
    pic.left = x + int((w - pic.width) / 2)
    pic.top = y + int((h - pic.height) / 2)
    return pic


def add_callouts(slide, callouts: list[tuple[str, str, str]], x=Inches(10.15), y=Inches(1.2), w=Inches(2.68)):
    for idx, (label, value, accent) in enumerate(callouts):
        top = y + Inches(idx * 1.02)
        add_metric_card(slide, x, top, w, Inches(0.86), label, value, accent=accent)


def add_figure_slide(prs, slide_no: int, fig_row: dict[str, str], section: str, title: str, callouts, note, bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, section, title)
    add_picture_fit(slide, Path(fig_row["png"]), Inches(0.42), Inches(1.12), Inches(9.55), Inches(5.72))
    add_callouts(slide, callouts)
    if bullets:
        add_multiline(slide, Inches(10.12), Inches(4.42), Inches(2.8), Inches(1.85), bullets, size=10.8)
    add_footer(slide, slide_no, note)
    return slide


def add_text_panel(slide, x, y, w, h, title, lines, accent="teal"):
    add_rect(slide, x, y, w, h, COLORS["white"], COLORS["line"], radius=True)
    add_rect(slide, x, y, Inches(0.07), h, COLORS[accent])
    add_text(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.28), Inches(0.28), title, 11, COLORS[accent], True)
    add_multiline(slide, x + Inches(0.16), y + Inches(0.52), w - Inches(0.28), h - Inches(0.62), lines, size=10.4)


def build_deck() -> tuple[Path, list[dict[str, str]], dict[str, float | int | str]]:
    for directory in [PPTX_DIR, REPORT_DIR, LOG_DIR]:
        directory.mkdir(parents=True, exist_ok=True)

    plot = load_plot_index()
    metrics = load_key_metrics()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slides: list[dict[str, str]] = []

    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["white"])
    add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, COLORS["navy"])
    add_text(slide, Inches(0.62), Inches(0.54), Inches(3.4), Inches(0.28), "GEO-ring Cloud / Stage 10", 11, COLORS["teal"], True, font=FONT_EN)
    add_text(slide, Inches(0.62), Inches(1.05), Inches(9.7), Inches(0.9), "GEO-ring 融合 CTH 产品验证与机制诊断", 31, COLORS["ink"], True)
    add_text(slide, Inches(0.66), Inches(2.12), Inches(9.9), Inches(0.45), "组会进度汇报：从 Stage09F 空间诊断到 Stage10 CTH 误差闭环", 17, COLORS["muted"])
    add_metric_card(slide, Inches(0.72), Inches(3.18), Inches(2.42), Inches(1.14), "D1 both-cloud n", fmt_num(metrics["d1_n"]), "像元级有效 CTH 对", "teal")
    add_metric_card(slide, Inches(3.38), Inches(3.18), Inches(2.25), Inches(1.14), "MAE", f"{metrics['d1_mae']:.3f} km", "融合对 EPIC A-band", "red")
    add_metric_card(slide, Inches(5.86), Inches(3.18), Inches(2.25), Inches(1.14), "Regret", f"{metrics['regret_delta']:.3f} km", "相对 oracle 诊断", "gold")
    add_metric_card(slide, Inches(8.35), Inches(3.18), Inches(2.38), Inches(1.14), "FOV sensitivity", fmt_delta(metrics["box7_cth_mae_delta"], " km"), "box 7x7 vs nearest", "green")
    add_text_panel(
        slide,
        Inches(0.72),
        Inches(4.75),
        Inches(5.5),
        Inches(1.45),
        "一句话结论",
        [
            "融合 CTH 与 EPIC effective-height reference 存在系统正偏。",
            "误差不是单一边界效应，选源机制和高度语义共同参与。",
            "近似 EPIC-FOV 聚合能改善指标，但幅度不足以推翻主结论。",
        ],
        "navy",
    )
    add_text(slide, Inches(9.9), Inches(6.62), Inches(2.2), Inches(0.25), datetime.now().strftime("%Y-%m-%d"), 9, COLORS["muted"], False, PP_ALIGN.RIGHT, FONT_EN)
    add_footer(slide, 1, "All slides are generated from existing Stage 10 CSV/JSON/Figure artifacts.")
    slides.append({"slide": "1", "title": "封面", "evidence": "stage_10 meeting figure package"})

    slide = prs.slides.add_slide(blank)
    add_header(slide, "STORY", "这次汇报的证据链")
    x_positions = [Inches(0.65), Inches(3.72), Inches(6.79), Inches(9.86)]
    steps = [
        ("Stage09F", "cloud-mask 空间故事", ["agreement", "boundary/broken", "valid-source-count"]),
        ("Stage10", "融合 CTH 验证", ["bias / MAE / RMSE", "within-2km", "D0-D7 domains"]),
        ("Mechanism", "误差来源诊断", ["selected source", "prefusion source", "selection regret"]),
        ("Stage10P/P2", "Composite 与 FOV 敏感性", ["无显式 kernel", "official PSF-aware", "approx FOV pilot"]),
    ]
    for idx, (tag, title, lines) in enumerate(steps):
        add_rect(slide, x_positions[idx], Inches(1.78), Inches(2.38), Inches(2.42), COLORS["white"], COLORS["line"], radius=True)
        add_text(slide, x_positions[idx] + Inches(0.18), Inches(1.98), Inches(2.0), Inches(0.28), tag, 13, COLORS["teal"], True, font=FONT_EN)
        add_text(slide, x_positions[idx] + Inches(0.18), Inches(2.38), Inches(2.0), Inches(0.34), title, 15, COLORS["ink"], True)
        add_multiline(slide, x_positions[idx] + Inches(0.22), Inches(2.93), Inches(1.95), Inches(0.85), lines, 10.3, COLORS["muted"])
        if idx < 3:
            add_text(slide, x_positions[idx] + Inches(2.51), Inches(2.76), Inches(0.42), Inches(0.28), ">", 21, COLORS["gold"], True, PP_ALIGN.CENTER, FONT_EN)
    add_text_panel(
        slide,
        Inches(0.75),
        Inches(4.78),
        Inches(5.55),
        Inches(1.25),
        "汇报主轴",
        ["先证明为什么要做 CTH，再展示主误差，再拆机制，最后用 Composite/FOV 把解释闭环。"],
        "navy",
    )
    add_text_panel(
        slide,
        Inches(6.75),
        Inches(4.78),
        Inches(5.55),
        Inches(1.25),
        "边界口径",
        ["EPIC 是 A-band Effective Cloud Height reference；Composite 可作为 PSF-aware benchmark 线索，不声称有官方 PSF kernel。"],
        "red",
    )
    add_footer(slide, 2)
    slides.append({"slide": "2", "title": "证据链", "evidence": "figure guide + source data"})

    slide = prs.slides.add_slide(blank)
    add_header(slide, "RESULT", "Stage10 fused CTH 的核心数值")
    cards = [
        ("n", fmt_num(metrics["d1_n"]), "D1 both-cloud 有效像元", "teal"),
        ("bias", f"{metrics['d1_bias']:+.3f} km", "GEO-ring 高于 reference", "red"),
        ("MAE", f"{metrics['d1_mae']:.3f} km", "平均绝对误差", "red"),
        ("RMSE", f"{metrics['d1_rmse']:.3f} km", "大误差更敏感", "gold"),
        ("within 2 km", f"{metrics['d1_within2']:.3f}", "2 km 内比例", "green"),
        ("high cloud MAE", f"{metrics['d7_mae']:.3f} km", "D7 高云域", "red"),
    ]
    for idx, card in enumerate(cards):
        col = idx % 3
        row = idx // 3
        add_metric_card(slide, Inches(0.72 + col * 3.95), Inches(1.35 + row * 1.42), Inches(3.25), Inches(1.1), *card)
    add_text_panel(
        slide,
        Inches(0.72),
        Inches(4.65),
        Inches(5.62),
        Inches(1.35),
        "样本与结果",
        [
            "D1 both-cloud 包含 3362.8 万个有效 CTH 像元。",
            "正 bias 与高云高 MAE 同时出现，指向高度语义和源选择机制。",
        ],
        "navy",
    )
    add_text_panel(
        slide,
        Inches(6.75),
        Inches(4.65),
        Inches(5.62),
        Inches(1.35),
        "解释边界",
        [
            "EPIC A-band 是 effective-height reference。",
            "2024-01 Composite 与 2024-03 Stage10 不做跨月数值对比。",
        ],
        "red",
    )
    add_footer(slide, 3)
    slides.append({"slide": "3", "title": "核心数值", "evidence": str(SOURCE_DIR / "stage_10_group_meeting_fig02_source.csv")})

    add_figure_slide(
        prs,
        4,
        plot["fig01"],
        "FIG. 1",
        "Stage09F 到 Stage10：为什么要从 mask 走到 CTH",
        [
            ("Stage09F", "空间诊断", "teal"),
            ("Stage10", "高度验证", "navy"),
            ("Logic", "问题递进", "gold"),
        ],
        "Stage09F is input evidence; Stage10 CTH validation is the current reporting target.",
        ["看点：agreement 不是终点；边界、碎云和选源数量把 CTH 误差问题推出来。"],
    )
    slides.append({"slide": "4", "title": "Fig. 1 证据链", "evidence": plot["fig01"]["source_csv"]})

    add_figure_slide(
        prs,
        5,
        plot["fig02"],
        "FIG. 2",
        "融合 CTH 主结果：D0-D7 诊断域",
        [
            ("D1 n", fmt_num(metrics["d1_n"]), "teal"),
            ("D1 MAE", f"{metrics['d1_mae']:.3f} km", "red"),
            ("within 2 km", f"{metrics['d1_within2']:.3f}", "green"),
        ],
        "Domain metrics are EPIC A-band Effective Cloud Height referenced.",
        ["D4 为 0 是当前策略与样本定义下的诊断结果，不单独解释成物理缺失。"],
    )
    slides.append({"slide": "5", "title": "Fig. 2 主结果", "evidence": plot["fig02"]["source_csv"]})

    add_figure_slide(
        prs,
        6,
        plot["fig03"],
        "FIG. 3",
        "误差来源分解：selected source 与 prefusion source",
        [
            ("Question", "源主导？", "gold"),
            ("Selected", "生产选择", "navy"),
            ("Prefusion", "可用源池", "teal"),
        ],
        "Source decomposition compares diagnostic source groups under the same EPIC reference.",
        ["如果被问“是不是某颗卫星拖累”，回答要回到 selected 与 prefusion 的并列证据。"],
    )
    slides.append({"slide": "6", "title": "Fig. 3 源分解", "evidence": plot["fig03"]["source_csv"]})

    add_figure_slide(
        prs,
        7,
        plot["fig04"],
        "FIG. 4",
        "语义审计与 A/B-band 敏感性",
        [
            ("EPIC", "effective height", "red"),
            ("D1 B-A", fmt_delta(metrics["ab_d1_mae_delta"], " km"), "gold"),
            ("Use", "reference", "teal"),
        ],
        "EPIC A/B-band heights are effective-height references; GEO product height is CTH.",
        ["A/B-band 差异用于量化 reference 语义敏感性，而不是替换主验证基准。"],
    )
    slides.append({"slide": "7", "title": "Fig. 4 语义审计", "evidence": plot["fig04"]["source_csv"]})

    add_figure_slide(
        prs,
        8,
        plot["fig05"],
        "FIG. 5",
        "selection regret 与高云/边界机制诊断",
        [
            ("current MAE", f"{metrics['regret_current']:.3f} km", "red"),
            ("oracle MAE", f"{metrics['regret_oracle']:.3f} km", "green"),
            ("regret", f"{metrics['regret_delta']:.3f} km", "gold"),
        ],
        "Oracle is a retrospective EPIC-referenced diagnostic, not a deployable production rule.",
        [
            f"高云 regret = {metrics['regret_high']:.3f} km。",
            f"clean-core MAE {metrics['clean_mae']:.3f} km，高于 boundary/broken {metrics['boundary_mae']:.3f} km。",
        ],
    )
    slides.append({"slide": "8", "title": "Fig. 5 机制诊断", "evidence": plot["fig05"]["source_csv"]})

    add_figure_slide(
        prs,
        9,
        plot["fig06"],
        "FIG. 6",
        "Stage10P/10P2：Composite 与近似 FOV 机制闭环",
        [
            ("Composite OK", f"{metrics['composite_files']}", "teal"),
            ("PSF kernel", f"{metrics['explicit_psf_kernel']}", "red"),
            ("box 7x7 MAE", fmt_delta(metrics["box7_cth_mae_delta"], " km"), "green"),
        ],
        "Composite is treated as official PSF-aware benchmark evidence, not as official PSF kernel.",
        [
            f"显式 weight candidate = {metrics['explicit_weight']}。",
            f"高云 MAE delta = {metrics['box7_high_mae_delta']:.3f} km；mask agreement delta = {metrics['box7_agree_delta']:+.4f}。",
        ],
    )
    slides.append({"slide": "9", "title": "Fig. 6 Composite/FOV", "evidence": plot["fig06"]["source_csv"]})

    slide = prs.slides.add_slide(blank)
    add_header(slide, "INTERPRETATION", "结论、证据与解释边界")
    rows = [
        ("融合 CTH 整体偏高", f"D1 bias {metrics['d1_bias']:+.3f} km；MAE {metrics['d1_mae']:.3f} km。", "相对 EPIC A-band ECH reference"),
        ("高云是关键误差域", f"D7 MAE {metrics['d7_mae']:.3f} km；high-cloud regret {metrics['regret_high']:.3f} km。", "需继续按 phase/COT/height 分层"),
        ("选源机制有改进空间", f"current MAE {metrics['regret_current']:.3f} km；oracle {metrics['regret_oracle']:.3f} km。", "oracle 是回溯诊断，不是生产规则"),
        ("近似 FOV 改善有限", f"box 7x7 D1 MAE delta {metrics['box7_cth_mae_delta']:.3f} km。", "无官方 PSF kernel 数值"),
    ]
    widths = [Inches(2.4), Inches(6.8), Inches(2.4)]
    x0 = Inches(0.65)
    y0 = Inches(1.36)
    add_rect(slide, x0, y0, sum(widths, Inches(0)), Inches(0.36), COLORS["navy"])
    headers = ["科学结论", "数据证据", "解释边界"]
    x = x0
    for header, width in zip(headers, widths):
        add_text(slide, x + Inches(0.1), y0 + Inches(0.07), width - Inches(0.2), Inches(0.18), header, 9.5, COLORS["white"], True)
        x += width
    for idx, row in enumerate(rows):
        y = y0 + Inches(0.42 + idx * 0.88)
        add_rect(slide, x0, y, sum(widths, Inches(0)), Inches(0.78), COLORS["white"], COLORS["line"])
        x = x0
        for col, width in zip(row, widths):
            add_text(slide, x + Inches(0.1), y + Inches(0.10), width - Inches(0.2), Inches(0.48), col, 10.1, COLORS["ink"])
            x += width
    add_text_panel(
        slide,
        Inches(0.72),
        Inches(5.55),
        Inches(11.6),
        Inches(0.8),
        "汇报定位",
        ["本报告给出 GEO-ring fused CTH 相对 EPIC effective-height reference 的一致性评估，并用源分解、regret 与 FOV 敏感性解释误差机制。"],
        "gold",
    )
    add_footer(slide, 10)
    slides.append({"slide": "10", "title": "结论与解释边界", "evidence": str(INTERROGATION_GUIDE)})

    slide = prs.slides.add_slide(blank)
    add_header(slide, "NEXT", "后续工作：可验证的改进路径")
    next_steps = [
        ("2024-01 pilot", ["跑一小批 GEO-ring，与官方 Composite 做同月 benchmark。", "只比较同月、同 FOV 语义下的结果。"], "teal"),
        ("Source-aware correction", ["针对 selected/prefusion source 分解，先做源别偏差校正。", "优先检查 Meteosat 与高云域。"], "navy"),
        ("High-cloud mechanism", ["把 D7 高云从单一阈值扩展为 phase/optical/height 分层。", "避免把所有高云误差混成一个机制。"], "red"),
        ("Governance package", ["保留 source data、manifest、QA 与脚本。", "后续图表可一键重建。"], "green"),
    ]
    for idx, (title, lines, accent) in enumerate(next_steps):
        col = idx % 2
        row = idx // 2
        add_text_panel(
            slide,
            Inches(0.75 + col * 6.0),
            Inches(1.35 + row * 2.0),
            Inches(5.2),
            Inches(1.42),
            title,
            lines,
            accent,
        )
    add_footer(slide, 11, "Next actions keep month, reference semantics, and product lineage explicit.")
    slides.append({"slide": "11", "title": "下一步", "evidence": "stage_10 group meeting package"})

    slide = prs.slides.add_slide(blank)
    add_header(slide, "APPENDIX", "变量定义与单位")
    glossary = [
        ("CTH", "Cloud Top Height；GEO 产品中的云顶高度变量，单位 km。"),
        ("EPIC ECH", "Effective Cloud Height；A/B-band 氧气吸收带反演的有效云高 reference。"),
        ("D1", "both-cloud；EPIC 与 GEO 都判云且有有效 CTH 的核心验证域。"),
        ("MAE/RMSE", "平均绝对误差/均方根误差；RMSE 对极端误差更敏感。"),
        ("within 2 km", "绝对误差小于等于 2 km 的像元比例。"),
        ("regret", "当前选源 MAE 与同像元 best-available prefusion source MAE 的差。"),
        ("PSF-aware", "Composite 结果含 FOV/PSF 加权证据；不等于文件中有 kernel 数值。"),
        ("box 7x7", "Stage10P2 近似 EPIC-FOV 聚合方法之一，不是官方 PSF。"),
    ]
    for idx, (term, definition) in enumerate(glossary):
        col = idx % 2
        row = idx // 2
        x = Inches(0.75 + col * 6.0)
        y = Inches(1.22 + row * 1.13)
        add_rect(slide, x, y, Inches(5.25), Inches(0.82), COLORS["white"], COLORS["line"], radius=True)
        add_text(slide, x + Inches(0.18), y + Inches(0.14), Inches(1.25), Inches(0.22), term, 11.2, COLORS["teal"], True, font=FONT_EN if re.match(r"^[A-Za-z0-9 /]+$", term) else FONT_CN)
        add_text(slide, x + Inches(1.48), y + Inches(0.10), Inches(3.55), Inches(0.44), definition, 9.2, COLORS["ink"])
    add_footer(slide, 12)
    slides.append({"slide": "12", "title": "变量定义与单位", "evidence": str(FIGURE_GUIDE)})

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX, slides, metrics


def write_slide_index(slides: list[dict[str, str]]) -> None:
    with SLIDE_INDEX.open("w", encoding="utf-8-sig", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=["slide", "title", "evidence"])
        writer.writeheader()
        writer.writerows(slides)


def write_speaker_notes(slides: list[dict[str, str]], metrics: dict[str, float | int | str]) -> None:
    lines = [
        "# Stage 10 组会 PPT 逐页讲稿备注",
        "",
        f"- PPTX: `{OUTPUT_PPTX}`",
        f"- 生成时间 UTC: `{datetime.now(timezone.utc).isoformat(timespec='seconds')}`",
        "- 证据边界：仅使用本地已有 CSV/JSON/图版；不联网、不重跑 Stage05/06；EPIC 写作 Effective Cloud Height reference。",
        "",
        "## Slide 1 封面",
        f"开场句：这次汇报聚焦 GEO-ring 融合 CTH 产品相对于 EPIC A-band Effective Cloud Height reference 的验证。D1 both-cloud 有 `{fmt_num(metrics['d1_n'])}` 个有效像元，MAE 为 `{metrics['d1_mae']:.3f} km`。",
        "",
        "## Slide 2 证据链",
        "讲法：先用 Stage09F 说明 cloud-mask 空间一致性并不能回答高度误差问题，再进入 Stage10 的 CTH 指标、源分解、selection regret，最后用 Stage10P/10P2 回答 EPIC FOV/PSF 机制。",
        "",
        "## Slide 3 核心数值",
        f"重点念：D1 n `{fmt_num(metrics['d1_n'])}`，bias `{metrics['d1_bias']:+.3f} km`，MAE `{metrics['d1_mae']:.3f} km`，RMSE `{metrics['d1_rmse']:.3f} km`，within-2km `{metrics['d1_within2']:.3f}`。这页只给结论，不展开机制。",
        "",
        "## Slide 4 Fig. 1",
        "讲法：Stage09F 的 agreement、boundary/broken、Meteosat-selected 和 valid-source-count 是空间证据链，说明下一步必须问“高度产品是否也一致”。",
        "",
        "## Slide 5 Fig. 2",
        "讲法：横向比较 D0-D7，各域同时看 bias、MAE、RMSE 和 within-2km。D7 high cloud MAE 最高，说明高云是 Stage10 误差诊断重点。",
        "",
        "## Slide 6 Fig. 3",
        "讲法：把生产 selected source 和 prefusion source pool 分开，是为了回答误差是否由某个 GEO 源或选源机制主导。不要只凭单颗卫星下结论，要看两组指标是否一致支持。",
        "",
        "## Slide 7 Fig. 4",
        f"讲法：EPIC A/B-band 是 effective-height reference，不是严格几何 CTH。D1 中 B-band 相对 A-band 的 MAE 增量为 `{metrics['ab_d1_mae_delta']:+.3f} km`，所以 reference 语义差异是必须报告的系统不确定性。",
        "",
        "## Slide 8 Fig. 5",
        f"讲法：当前 selected MAE `{metrics['regret_current']:.3f} km`，best-available oracle `{metrics['regret_oracle']:.3f} km`，regret `{metrics['regret_delta']:.3f} km`。高云 regret `{metrics['regret_high']:.3f} km`，提示选源机制在高云域有改进空间。",
        "",
        "## Slide 9 Fig. 6",
        f"讲法：Stage10P 打开 `{metrics['composite_files']}` 个 2024-01 Composite 文件，没有显式 PSF kernel/weight 数值，因此只能说 official PSF-aware benchmark。Stage10P2 中 box 7x7 对 D1 CTH MAE 的 delta 为 `{metrics['box7_cth_mae_delta']:.3f} km`，有改善但不是决定性主因。",
        "",
        "## Slide 10 导师追问口径",
        "讲法：这页用于答辩，不主动讲太久。核心是四条边界：EPIC 不是绝对真值；GOES HT 按 GEO CTH 产品变量处理；边界不是唯一解释；近似 FOV 改善有限。",
        "",
        "## Slide 11 下一步",
        "讲法：下一步不是盲目扩大数据，而是先做 2024-01 official Composite benchmark pilot，再针对源别偏差、高云域和 governance package 做可验证改进。",
        "",
        "## Slide 12 变量缩写",
        "讲法：这页是备份页。导师问变量定义时直接回到这里，尤其强调 CTH 与 EPIC Effective Cloud Height 的语义差异。",
        "",
    ]
    SPEAKER_NOTES.write_text("\n".join(lines), encoding="utf-8")


def inspect_pptx(path: Path) -> dict[str, object]:
    with zipfile.ZipFile(path, "r") as archive:
        names = archive.namelist()
        slide_xml = [name for name in names if re.match(r"ppt/slides/slide\d+\.xml$", name)]
        media = [name for name in names if name.startswith("ppt/media/")]
        has_content_types = "[Content_Types].xml" in names
        core_text = "\n".join(
            archive.read(name).decode("utf-8", errors="replace")
            for name in slide_xml[:2]
        )
    return {
        "pptx_exists": path.exists(),
        "pptx_size_bytes": path.stat().st_size if path.exists() else 0,
        "slide_count": len(slide_xml),
        "media_count": len(media),
        "has_content_types": has_content_types,
        "cover_text_present": "GEO-ring" in core_text and "Stage 10" in core_text,
    }


def scan_visible_pptx_text(path: Path) -> dict[str, object]:
    internal_terms = [
        "我会怎么说",
        "不能怎么说",
        "导师",
        "防守",
        "绝对真值",
        "not absolute",
        "truth",
    ]
    hits: list[dict[str, object]] = []
    prs = Presentation(str(path))
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs)
            for term in internal_terms:
                if term in text:
                    hits.append({"slide": slide_idx, "term": term, "text": text[:180]})
    return {"passed": len(hits) == 0, "hits": hits}


def scan_mojibake(paths: list[Path]) -> dict[str, object]:
    patterns = [
        "\u951f",
        "\ufffd",
        "\u00c3",
        "\u00c2",
        "\u00e4\u00bd",
        "\u00e6",
        "\u00ca",
        "\u00d0",
    ]
    hits: list[dict[str, str | int]] = []
    for path in paths:
        if not path.exists():
            hits.append({"path": str(path), "line": 0, "text": "missing"})
            continue
        for line_no, line in enumerate(path.read_text(encoding="utf-8", errors="replace").splitlines(), start=1):
            if any(pattern in line for pattern in patterns):
                hits.append({"path": str(path), "line": line_no, "text": line[:180]})
    return {"passed": len(hits) == 0, "hits": hits}


def inspect_render_check_pdf(path: Path) -> dict[str, object]:
    if not path.exists():
        return {"passed": False, "status": "not_run", "path": str(path)}
    try:
        from pypdf import PdfReader
    except ImportError:
        return {
            "passed": False,
            "status": "pypdf_unavailable",
            "path": str(path),
            "size_bytes": path.stat().st_size,
        }
    reader = PdfReader(str(path))
    sample_text = " ".join((reader.pages[idx].extract_text() or "") for idx in [0, min(4, len(reader.pages) - 1)])
    sample_text = sample_text.replace("\n", " ")
    latin_anchor_present = "GEO-ring" in sample_text and "33,627,889" in sample_text
    return {
        "passed": len(reader.pages) == 12 and path.stat().st_size > 100_000 and latin_anchor_present,
        "status": "ok",
        "path": str(path),
        "size_bytes": path.stat().st_size,
        "page_count": len(reader.pages),
        "latin_anchor_present": latin_anchor_present,
        "text_extraction_note": "PowerPoint exported the PDF successfully; CJK glyph fidelity is checked by Office rendering, not by pypdf text extraction.",
    }


def write_manifest_and_qa(pptx_path: Path, slides: list[dict[str, str]], metrics: dict[str, float | int | str]) -> None:
    manifest = {
        "stage_id": STAGE_ID,
        "run_id": RUN_ID,
        "deck_id": DECK_ID,
        "created_utc": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "script": str(SCRIPT_PATH),
        "pptx": str(pptx_path),
        "powerpoint_render_check_pdf": str(RENDER_CHECK_PDF),
        "slide_index": str(SLIDE_INDEX),
        "speaker_notes": str(SPEAKER_NOTES),
        "source_figure_index": str(PLOT_INDEX),
        "input_reports": [str(FIGURE_GUIDE), str(INTERROGATION_GUIDE)],
        "slide_count": len(slides),
        "key_metrics": metrics,
        "rules": [
            "No network download.",
            "No Stage05/06 rerun.",
            "EPIC is reported as A-band Effective Cloud Height reference.",
            "Composite is reported as official PSF-aware benchmark evidence unless explicit kernel/weight arrays are present.",
        ],
    }
    MANIFEST.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    pptx_info = inspect_pptx(pptx_path)
    visible_text = scan_visible_pptx_text(pptx_path)
    mojibake = scan_mojibake([SLIDE_INDEX, SPEAKER_NOTES, MANIFEST])
    qa = {
        "created_utc": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "checks": {
            "pptx_structure": {
                "passed": bool(
                    pptx_info["pptx_exists"]
                    and pptx_info["pptx_size_bytes"] > 100_000
                    and pptx_info["slide_count"] == 12
                    and pptx_info["media_count"] >= 6
                    and pptx_info["has_content_types"]
                ),
                "details": pptx_info,
            },
            "slide_index": {
                "passed": SLIDE_INDEX.exists() and len(read_csv_rows(SLIDE_INDEX)) == 12,
                "details": {"path": str(SLIDE_INDEX), "rows": len(read_csv_rows(SLIDE_INDEX)) if SLIDE_INDEX.exists() else 0},
            },
            "speaker_notes": {
                "passed": SPEAKER_NOTES.exists() and SPEAKER_NOTES.stat().st_size > 1000,
                "details": {"path": str(SPEAKER_NOTES), "bytes": SPEAKER_NOTES.stat().st_size if SPEAKER_NOTES.exists() else 0},
            },
            "mojibake_scan": mojibake,
            "powerpoint_pdf_render": inspect_render_check_pdf(RENDER_CHECK_PDF),
            "visible_pptx_text_audience_scan": visible_text,
        },
    }
    overall = all(check.get("passed", False) for check in qa["checks"].values())
    qa["overall_passed"] = overall
    QA_REPORT.write_text(json.dumps(qa, ensure_ascii=False, indent=2), encoding="utf-8")


def main(argv: list[str] | None = None) -> None:
    parser = build_parser()
    parser.parse_args(argv)
    missing = [path for path in [PLOT_INDEX, FIGURE_GUIDE, INTERROGATION_GUIDE] if not path.exists()]
    if missing:
        raise FileNotFoundError("Missing required Stage 10 meeting figure artifacts: " + ", ".join(map(str, missing)))
    pptx_path, slides, metrics = build_deck()
    write_slide_index(slides)
    write_speaker_notes(slides, metrics)
    write_manifest_and_qa(pptx_path, slides, metrics)
    print(json.dumps({"pptx": str(pptx_path), "slides": len(slides), "qa": str(QA_REPORT)}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
